import json
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def fetch_random_image():
    # Example of a random image API (Unsplash)
    url = "https://source.unsplash.com/random/800x600"
    response = requests.get(url)
    return BytesIO(response.content)

def create_presentation(json_file, output_pptx):
    prs = Presentation()

    with open(json_file, 'r') as f:
        data = json.load(f)

    for slide_data in data:
        slide_layout = prs.slide_layouts[5]  # Using blank slide layout

        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data.get('title', 'Title')

        # Adding image
        left_img = Inches(0.5)
        top_img = Inches(1.5)
        width_img = Inches(5.5)
        height_img = Inches(5.5)
        img_stream = fetch_random_image()
        pic = slide.shapes.add_picture(img_stream, left_img, top_img, width_img, height_img)

        # Adding text
        left_text = Inches(6.0)
        top_text = Inches(1.5)
        width_text = Inches(5.5)
        height_text = Inches(5.5)

        txBox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = slide_data.get('description_title', 'Description')
        p.font.bold = True
        p.font.size = Pt(20)

        for bullet in slide_data['description']:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save(output_pptx)
    print("Presentation created successfully!")

# Example usage:
json_file = 'presentation_dat.json'
with open('transcribed_text.txt','r') as tttt:
    otpt = tttt.read()
output_pptx = f'output/{otpt.strip()}.pptx'
create_presentation(json_file, output_pptx)